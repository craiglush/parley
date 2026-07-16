"""Attachment text extraction — turn any attachment into text (or defer to GPU).

extract_text(path, filename) -> Extraction dict {text, method, chars, status}
with status in {done, pending, empty, failed}.

INSTANT methods (plaintext incl. svg, PDF text layer, docx/xlsx/pptx) parse
synchronously here; the CALLER must run extract_text in a thread executor
(run_in_executor / _run_bg style) so a 50 MiB parse never blocks the event loop.
DEFERRED methods (audio/video STT, raster-image / scanned-PDF vision) return
status='pending' with the method to run later; a background worker resolves them
off the GPU-contended path.

Every extractor is wrapped so a parse error yields status='failed' and never
raises — attachment storage must never break.
"""
import io  # noqa: F401  (used by later tasks: render_pdf_page_pngs)
import json  # noqa: F401  (used by later tasks: sidecar read/write)
import logging
from datetime import datetime, timezone  # noqa: F401  (used by write_extraction)
from pathlib import Path

logger = logging.getLogger("meeting-service")

try:
    from pypdf import PdfReader
except Exception:  # dep not installed yet / import error -> pdf path degrades to failed
    PdfReader = None

# Plaintext families (incl. svg as XML text — never sent to the vision model).
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log", ".yaml", ".yml",
    ".xml", ".html", ".htm", ".svg", ".rst", ".ini", ".toml", ".env", ".conf", ".cfg",
    ".py", ".js", ".mjs", ".ts", ".tsx", ".jsx", ".css", ".scss", ".sh", ".bat",
    ".ps1", ".c", ".h", ".cpp", ".hpp", ".java", ".go", ".rs", ".rb", ".php", ".sql",
}
_AV_EXTS = {
    ".mp3", ".wav", ".m4a", ".flac", ".ogg", ".oga", ".opus", ".aac", ".wma",
    ".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v", ".wmv", ".flv",
}
_RASTER_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"}


def _ext_of(filename: str) -> str:
    return ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""


def _result(text: str, method: str, status: str) -> dict:
    text = text or ""
    return {"text": text, "method": method, "chars": len(text), "status": status}


def _extract_plaintext(path: str) -> dict:
    text = Path(path).read_bytes().decode("utf-8", errors="replace")
    return _result(text, "text", "done" if text.strip() else "empty")


def extract_text(path: str, filename: str) -> dict:
    """Dispatch by extension. Instant methods parse now; deferred return pending.
    Any extractor failure -> status='failed' (never raises)."""
    ext = _ext_of(filename)
    try:
        if ext in _TEXT_EXTS:
            return _extract_plaintext(path)
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".xlsx":
            return _extract_xlsx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext in _AV_EXTS:
            return _result("", "stt", "pending")
        if ext in _RASTER_EXTS:
            return _result("", "vision", "pending")
        return _result("", "", "empty")
    except Exception as e:
        logger.warning(f"extract_text failed for {filename!r} (non-fatal): {e}")
        return _result("", "", "failed")


EXTRACTED_DIRNAME = ".extracted"


def _now_iso() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def extracted_sidecar_path(attach_dir, stored_filename: str) -> Path:
    """Per-attachment sidecar path: <attach_dir>/.extracted/<stored_filename>.json."""
    return Path(attach_dir) / EXTRACTED_DIRNAME / f"{stored_filename}.json"


def write_extraction(attach_dir, stored_filename: str, extraction: dict) -> Path:
    """Atomically write the .extracted sidecar (extraction shape + extracted_at)."""
    p = extracted_sidecar_path(attach_dir, stored_filename)
    p.parent.mkdir(parents=True, exist_ok=True)
    record = {**extraction, "extracted_at": _now_iso()}
    tmp = p.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(record, ensure_ascii=False), encoding="utf-8")
    tmp.replace(p)
    return p


def read_extraction(attach_dir, stored_filename: str) -> dict | None:
    """Read a .extracted sidecar; None if absent or corrupt. (Phase B/C reader.)"""
    p = extracted_sidecar_path(attach_dir, stored_filename)
    if not p.exists():
        return None
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else None
    except (OSError, ValueError):
        return None


# Instant document extractors (deps land in Task 2). Kept above extract_text's
# reference via late binding — Python resolves them at call time.
_PDF_MIN_TEXT_CHARS = 16  # below this the text layer is "empty" -> scanned -> vision


def _extract_pdf(path: str) -> dict:
    if PdfReader is None:
        raise RuntimeError("pypdf not available")
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            parts.append(t)
    text = "\n".join(parts)
    if sum(1 for c in text if not c.isspace()) < _PDF_MIN_TEXT_CHARS:
        return _result("", "vision", "pending")  # scanned -> deferred vision (rasterize later)
    return _result(text, "pdf", "done")


def _extract_docx(path: str) -> dict:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    text = "\n".join(parts)
    return _result(text, "docx", "done" if text.strip() else "empty")


def _extract_xlsx(path: str) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
    finally:
        wb.close()
    text = "\n".join(parts)
    return _result(text, "xlsx", "done" if text.strip() else "empty")


def _extract_pptx(path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    parts.append(line)
    text = "\n".join(parts)
    return _result(text, "pptx", "done" if text.strip() else "empty")


def render_pdf_page_pngs(path: str, *, max_pages: int = 8, scale: float = 2.0) -> list:
    """Rasterize a (scanned) PDF's pages to PNG bytes for the vision model.
    Capped at max_pages so a huge scan can't fan out into hundreds of GPU calls."""
    import pypdfium2 as pdfium
    pdf = pdfium.PdfDocument(path)
    out = []
    try:
        n = min(len(pdf), max_pages)
        for i in range(n):
            page = pdf[i]
            bitmap = page.render(scale=scale)
            pil = bitmap.to_pil()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            out.append(buf.getvalue())
    finally:
        pdf.close()
    return out
