import extract


def _write(tmp_path, name, data=b"hello"):
    p = tmp_path / name
    p.write_bytes(data)
    return str(p)


def test_plaintext_txt_done(tmp_path):
    r = extract.extract_text(_write(tmp_path, "a.txt", b"hello world"), "a.txt")
    assert r == {"text": "hello world", "method": "text", "chars": 11, "status": "done"}


def test_csv_is_plaintext(tmp_path):
    r = extract.extract_text(_write(tmp_path, "a.csv", b"x,y\n1,2"), "a.csv")
    assert r["status"] == "done" and r["method"] == "text" and "1,2" in r["text"]


def test_svg_decoded_as_text_not_vision(tmp_path):
    svg = b"<svg><text>diagram label</text></svg>"
    r = extract.extract_text(_write(tmp_path, "d.svg", svg), "d.svg")
    assert r["method"] == "text" and r["status"] == "done"
    assert "diagram label" in r["text"]


def test_empty_plaintext_is_empty_status(tmp_path):
    r = extract.extract_text(_write(tmp_path, "blank.txt", b"   \n"), "blank.txt")
    assert r["status"] == "empty" and r["chars"] == 4


def test_audio_is_pending_stt(tmp_path):
    r = extract.extract_text(_write(tmp_path, "rec.m4a"), "rec.m4a")
    assert r == {"text": "", "method": "stt", "chars": 0, "status": "pending"}


def test_raster_image_is_pending_vision(tmp_path):
    r = extract.extract_text(_write(tmp_path, "pic.png"), "pic.png")
    assert r == {"text": "", "method": "vision", "chars": 0, "status": "pending"}


def test_unknown_binary_is_empty(tmp_path):
    r = extract.extract_text(_write(tmp_path, "thing.bin"), "thing.bin")
    assert r == {"text": "", "method": "", "chars": 0, "status": "empty"}


def test_missing_file_is_failed_not_raised(tmp_path):
    r = extract.extract_text(str(tmp_path / "nope.txt"), "nope.txt")
    assert r["status"] == "failed" and r["chars"] == 0


def test_extraction_deps_importable():
    import pypdf            # noqa: F401
    import docx             # noqa: F401  (python-docx)
    import openpyxl         # noqa: F401
    import pptx             # noqa: F401  (python-pptx)
    import pypdfium2        # noqa: F401
    assert extract.PdfReader is not None


def test_sidecar_roundtrip(tmp_path):
    attach = tmp_path / "attachments"
    ex = {"text": "hi", "method": "text", "chars": 2, "status": "done"}
    p = extract.write_extraction(attach, "notes-ab12cd.txt", ex)
    assert p == attach / ".extracted" / "notes-ab12cd.txt.json"
    assert p.exists()
    got = extract.read_extraction(attach, "notes-ab12cd.txt")
    assert got["text"] == "hi" and got["status"] == "done"
    assert "extracted_at" in got  # added at write time, not part of the 4-key shape


def test_read_missing_sidecar_is_none(tmp_path):
    assert extract.read_extraction(tmp_path / "attachments", "missing.txt") is None


def test_read_corrupt_sidecar_is_none(tmp_path):
    attach = tmp_path / "attachments"
    p = extract.extracted_sidecar_path(attach, "x.txt")
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text("{not json", encoding="utf-8")
    assert extract.read_extraction(attach, "x.txt") is None


def test_read_extraction_invalid_utf8_returns_none(tmp_path):
    p = extract.extracted_sidecar_path(tmp_path, "x.txt")
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_bytes(b"\xff\xfe\x00garbage")
    assert extract.read_extraction(tmp_path, "x.txt") is None


def test_read_extraction_non_dict_json_returns_none(tmp_path):
    p = extract.extracted_sidecar_path(tmp_path, "y.txt")
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text("[]", encoding="utf-8")
    assert extract.read_extraction(tmp_path, "y.txt") is None


def test_pdf_text_layer_done(tmp_path, monkeypatch):
    class _Page:
        def extract_text(self):
            return "Quarterly revenue grew twelve percent this period."

    class _Reader:
        def __init__(self, path):
            self.pages = [_Page(), _Page()]

    monkeypatch.setattr(extract, "PdfReader", _Reader)
    p = tmp_path / "doc.pdf"
    p.write_bytes(b"%PDF-1.4 stub")
    r = extract.extract_text(str(p), "doc.pdf")
    assert r["method"] == "pdf" and r["status"] == "done"
    assert "Quarterly revenue" in r["text"]


def test_scanned_pdf_is_pending_vision(tmp_path):
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=72, height=72)  # real PDF, empty text layer
    p = tmp_path / "scan.pdf"
    with open(p, "wb") as fh:
        w.write(fh)
    r = extract.extract_text(str(p), "scan.pdf")
    assert r == {"text": "", "method": "vision", "chars": 0, "status": "pending"}


def test_pdf_parse_error_is_failed(tmp_path):
    p = tmp_path / "bad.pdf"
    p.write_bytes(b"not really a pdf")
    r = extract.extract_text(str(p), "bad.pdf")
    assert r["status"] in ("failed", "pending")  # unreadable -> failed; blank text -> pending


def test_pdf_sparse_whitespace_text_is_pending_vision(tmp_path, monkeypatch):
    class _Page:
        def extract_text(self):
            return "a" + " " * 30 + "b"

    class _Reader:
        def __init__(self, path):
            self.pages = [_Page()]

    monkeypatch.setattr(extract, "PdfReader", _Reader)
    p = tmp_path / "sparse.pdf"
    p.write_bytes(b"%PDF-fake")
    r = extract.extract_text(str(p), "sparse.pdf")
    assert r["status"] == "pending" and r["method"] == "vision"


def test_docx_extract(tmp_path):
    from docx import Document
    d = Document()
    d.add_paragraph("Meeting agenda for Monday.")
    d.add_paragraph("Discuss the budget.")
    p = tmp_path / "agenda.docx"
    d.save(str(p))
    r = extract.extract_text(str(p), "agenda.docx")
    assert r["method"] == "docx" and r["status"] == "done"
    assert "Meeting agenda" in r["text"] and "budget" in r["text"]


def test_xlsx_extract(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Region"
    ws["B1"] = "Sales"
    ws["A2"] = "EMEA"
    ws["B2"] = 4200
    p = tmp_path / "sheet.xlsx"
    wb.save(str(p))
    r = extract.extract_text(str(p), "sheet.xlsx")
    assert r["method"] == "xlsx" and r["status"] == "done"
    assert "Region" in r["text"] and "EMEA" in r["text"] and "4200" in r["text"]


def test_pptx_extract(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish w/ title
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Roadmap milestones for Q3"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    r = extract.extract_text(str(p), "deck.pptx")
    assert r["method"] == "pptx" and r["status"] == "done"
    assert "Roadmap milestones" in r["text"]


def test_render_pdf_page_pngs(tmp_path):
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=144, height=144)
    w.add_blank_page(width=144, height=144)
    p = tmp_path / "scan.pdf"
    with open(p, "wb") as fh:
        w.write(fh)
    pngs = extract.render_pdf_page_pngs(str(p))
    assert len(pngs) == 2
    assert all(b.startswith(b"\x89PNG") for b in pngs)


def test_render_pdf_page_pngs_caps_pages(tmp_path):
    from pypdf import PdfWriter
    w = PdfWriter()
    for _ in range(5):
        w.add_blank_page(width=72, height=72)
    p = tmp_path / "many.pdf"
    with open(p, "wb") as fh:
        w.write(fh)
    assert len(extract.render_pdf_page_pngs(str(p), max_pages=3)) == 3
